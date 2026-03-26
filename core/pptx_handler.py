from pptx import Presentation

class PptxHandler:
    @staticmethod
    def translate_pptx(file_path, translator, output_path=None):
        prs = Presentation(file_path)

        # 1. Collect all runs with text
        targets = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            targets.append(run)
        
        # 2. Extract text and batch translate
        texts_to_translate = [r.text for r in targets]
        translated_texts = translator.translate_texts(texts_to_translate)
        
        # 3. Apply back
        for run, translated_text in zip(targets, translated_texts):
            run.text = translated_text

        if output_path is None:
            output_path = file_path.replace(".pptx", "_translated.pptx")
        
        prs.save(output_path)
        return output_path
